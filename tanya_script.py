#!/usr/bin/env python3
"""
Sports Fan Sponsorship Analysis - Python Implementation
Updated with TRIM functions to handle whitespace in data
Fixed SQL escaping for merchant names with apostrophes
Fixed f-string syntax errors
Added Red Hat Display font support for PowerPoint
"""

import os
import pandas as pd
import numpy as np
from datetime import datetime
from pathlib import Path
import sys

# Add project root to path to import snowflake_connection
project_root = Path(__file__).resolve().parent
sys.path.append(str(project_root))

# Import your centralized connection manager
from snowflake_connection import get_connection, query_to_dataframe, with_snowflake_connection

# PowerPoint imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.table import Table

import warnings

warnings.filterwarnings('ignore')


class SponsorshipAnalyzer:
    """Main class for fan sponsorship analysis and reporting"""

    def __init__(self):
        """Initialize settings - connection handled by manager"""
        self.categories = ["Athleisure", "Restaurants", "Finance", "Auto", "Gambling", "Travel"]
        self.current_category = None

        # Setup Red Hat Display font
        self.setup_font()

    def setup_font(self):
        """Setup Red Hat Display font for PowerPoint"""
        try:
            # Get the script's directory
            script_dir = Path(__file__).parent

            # Check parent directories too (in case we're in graphics/slide 2)
            parent_dir = script_dir.parent
            grandparent_dir = parent_dir.parent

            # Define all possible font locations
            font_search_paths = [
                # From static directory in script directory
                script_dir / "static" / "RedHatDisplay-Bold.ttf",
                script_dir / "static" / "RedHatDisplay-Regular.ttf",
                script_dir / "static" / "RedHatDisplay-Medium.ttf",
                script_dir / "static" / "RedHatDisplay-BoldItalic.ttf",
                script_dir / "static" / "RedHatDisplay-Italic.ttf",
                # From parent directory static folder
                parent_dir / "static" / "RedHatDisplay-Bold.ttf",
                parent_dir / "static" / "RedHatDisplay-Regular.ttf",
                # From grandparent (project root) static folder
                grandparent_dir / "static" / "RedHatDisplay-Bold.ttf",
                grandparent_dir / "static" / "RedHatDisplay-Regular.ttf",
                # Also check the old paths in case they exist elsewhere
                script_dir / "Red_Hat_Display" / "static" / "RedHatDisplay-Bold.ttf",
                script_dir / "Red_Hat_Display" / "static" / "RedHatDisplay-Regular.ttf",
                # Direct file in various locations
                script_dir / "RedHatDisplay-Bold.ttf",
                script_dir / "RedHatDisplay-Regular.ttf",
                # From current working directory static folder
                Path.cwd() / "static" / "RedHatDisplay-Bold.ttf",
                Path.cwd() / "static" / "RedHatDisplay-Regular.ttf",
            ]

            # Track what we find
            self.font_paths = {
                'regular': None,
                'bold': None
            }

            # Search for fonts
            for font_path in font_search_paths:
                if font_path.exists():
                    if 'Bold' in font_path.name:
                        self.font_paths['bold'] = str(font_path)
                        print(f"✓ Found Bold font: {font_path.name}")
                    elif 'Regular' in font_path.name:
                        self.font_paths['regular'] = str(font_path)
                        print(f"✓ Found Regular font: {font_path.name}")

            if self.font_paths['regular'] or self.font_paths['bold']:
                self.font_name = 'Red Hat Display'
                print(f"✅ Red Hat Display font configured for PowerPoint")
            else:
                print("⚠️  Red Hat Display font not found - will use default fonts")
                self.font_name = None

        except Exception as e:
            print(f"❌ Error setting up font: {e}")
            self.font_name = None

    def _apply_font_to_text(self, paragraph, font_type='regular'):
        """Apply Red Hat Display font to a paragraph if available"""
        if self.font_name:
            # Try different font name variations that PowerPoint might recognize
            font_variations = [
                'Red Hat Display',  # Base name
                'RedHatDisplay',  # No spaces
                'Red Hat Display Regular',  # With style
                'RedHatDisplay-Regular'  # Hyphenated
            ]

            # Set the font name
            paragraph.font.name = self.font_name

            # Also try to ensure the font is applied at the run level
            if hasattr(paragraph, 'runs'):
                for run in paragraph.runs:
                    run.font.name = self.font_name

    def check_available_data(self):
        """Check what data is available in the tables"""
        print("\n🔍 Checking available data in Snowflake tables...")

        # Check categories in category table
        query = """
        SELECT DISTINCT TRIM(CATEGORY) as CATEGORY, COUNT(*) as ROW_COUNT
        FROM V_UTAH_JAZZ_SIL_CATEGORY_INDEXING_ALL_TIME
        WHERE AUDIENCE = 'Utah Jazz Fans'
        GROUP BY TRIM(CATEGORY)
        ORDER BY CATEGORY
        """
        categories_df = query_to_dataframe(query)
        print("\n📂 Available categories in CATEGORY table:")
        category_list = []
        for cat in categories_df['CATEGORY'].tolist():
            print(f"  - {cat}")
            category_list.append(cat)

        # Check merchant categories and subcategories
        query = """
        SELECT DISTINCT TRIM(CATEGORY) as CATEGORY, TRIM(SUBCATEGORY) as SUBCATEGORY, 
               COUNT(DISTINCT MERCHANT) as MERCHANT_COUNT
        FROM V_UTAH_JAZZ_SIL_MERCHANT_INDEXING_ALL_TIME
        WHERE AUDIENCE = 'Utah Jazz Fans'
        GROUP BY TRIM(CATEGORY), TRIM(SUBCATEGORY)
        ORDER BY CATEGORY, SUBCATEGORY
        """
        merchant_cats_df = query_to_dataframe(query)
        print("\n🏪 Categories and Subcategories in MERCHANT table:")
        for _, row in merchant_cats_df.iterrows():
            print(f"  - {row['CATEGORY']} > {row['SUBCATEGORY']}: {row['MERCHANT_COUNT']} merchants")

        return category_list

    def analyze_category(self, category):
        """Analyze a specific category using centralized connection"""
        self.current_category = category
        print(f"\n📊 Analyzing category: {category}")

        # 1. Category-level analysis
        category_stats = self._get_category_stats()

        # 2. Subcategory analysis
        subcategory_stats = self._get_subcategory_stats()

        # 3. Generate insights
        insights = self._generate_insights(category_stats, subcategory_stats)

        # 4. Merchant analysis
        merchant_stats = self._get_merchant_stats()

        # 5. Generate merchant insights
        merchant_insights = self._generate_merchant_insights(merchant_stats)

        return {
            'category_stats': category_stats,
            'subcategory_stats': subcategory_stats,
            'insights': insights,
            'merchant_stats': merchant_stats,
            'merchant_insights': merchant_insights
        }

    def _get_category_stats(self):
        """Get category-level statistics using query_to_dataframe"""
        query = f"""
        SELECT * FROM V_UTAH_JAZZ_SIL_CATEGORY_INDEXING_ALL_TIME 
        WHERE TRIM(CATEGORY) = %(category)s
        """

        # Use the centralized query function with parameters
        df = query_to_dataframe(query, {'category': self.current_category.strip()})

        if df.empty:
            print(f"⚠️  No category data found for: {self.current_category}")
            return {
                'percent_fans': "N/A",
                'likelihood': "N/A",
                'purchases': "N/A",
                'raw_values': {
                    'percent_fans': 0,
                    'percent_likely': 0,
                    'percent_purch': 0
                }
            }

        # Filter for Jazz fans
        all_fans = df[df['AUDIENCE'] == 'Utah Jazz Fans']

        if all_fans.empty:
            print(f"⚠️  No Utah Jazz Fans data found for category: {self.current_category}")
            return {
                'percent_fans': "N/A",
                'likelihood': "N/A",
                'purchases': "N/A",
                'raw_values': {
                    'percent_fans': 0,
                    'percent_likely': 0,
                    'percent_purch': 0
                }
            }

        percent_fans = float(all_fans['PERC_AUDIENCE'].iloc[0])

        # Compare to local general population
        comp_gen_pop = df[
            (df['AUDIENCE'] == 'Utah Jazz Fans') &
            (df['COMPARISON_POPULATION'] == 'Local Gen Pop (Excl. Jazz)')
            ]

        if comp_gen_pop.empty:
            print(f"⚠️  No comparison data found for Local Gen Pop")
            return {
                'percent_fans': f"{percent_fans * 100:.1f}%",
                'likelihood': "N/A",
                'purchases': "N/A",
                'raw_values': {
                    'percent_fans': percent_fans,
                    'percent_likely': 0,
                    'percent_purch': 0
                }
            }

        # Calculate metrics
        percent_likely = float(comp_gen_pop['PERC_INDEX'].iloc[0]) - 100
        percent_purch = self._calculate_percent_diff(
            float(comp_gen_pop['PPC'].iloc[0]),
            float(comp_gen_pop['COMPARISON_PPC'].iloc[0])
        )

        # Format results
        results = {
            'percent_fans': f"{percent_fans * 100:.1f}%",
            'likelihood': f"{abs(percent_likely):.1f}% {'More' if percent_likely > 0 else 'Less'}",
            'purchases': f"{abs(percent_purch):.1f}% {'More' if percent_purch > 0 else 'Less'}",
            'raw_values': {
                'percent_fans': percent_fans,
                'percent_likely': percent_likely,
                'percent_purch': percent_purch
            }
        }

        return results

    def _get_subcategory_stats(self):
        """Get subcategory-level statistics"""
        query = f"""
        SELECT * FROM V_UTAH_JAZZ_SIL_SUBCATEGORY_INDEXING_ALL_TIME 
        WHERE TRIM(CATEGORY) = %(category)s
        """

        df = query_to_dataframe(query, {'category': self.current_category.strip()})

        if df.empty:
            print(f"⚠️  No subcategory data found for: {self.current_category}")
            return pd.DataFrame()

        # Get top 5 subcategories
        all_fans = df[df['AUDIENCE'] == 'Utah Jazz Fans'].copy()
        if all_fans.empty:
            print(f"⚠️  No Utah Jazz Fans data in subcategories for: {self.current_category}")
            return pd.DataFrame()

        all_fans['PERC_AUDIENCE'] = pd.to_numeric(all_fans['PERC_AUDIENCE'])

        top_5_subs = (all_fans
                      .sort_values('PERC_AUDIENCE', ascending=False)
                      .drop_duplicates('SUBCATEGORY')
                      .head(5))

        results = []
        for _, row in top_5_subs.iterrows():
            subcategory = row['SUBCATEGORY']

            # Get comparison data
            comp_gen_pop = df[
                (df['AUDIENCE'] == 'Utah Jazz Fans') &
                (df['COMPARISON_POPULATION'] == 'Local Gen Pop (Excl. Jazz)') &
                (df['SUBCATEGORY'] == subcategory)
                ]

            if not comp_gen_pop.empty:
                percent_likely = float(comp_gen_pop['PERC_INDEX'].iloc[0]) - 100
                percent_purch = self._calculate_percent_diff(
                    float(comp_gen_pop['PPC'].iloc[0]),
                    float(comp_gen_pop['COMPARISON_PPC'].iloc[0])
                )

                results.append({
                    'Subcategory': subcategory,
                    'Percent of Fans Who Spend': f"{row['PERC_AUDIENCE'] * 100:.1f}%",
                    'How likely fans are to spend vs. gen pop': f"{abs(percent_likely):.1f}% {'More' if percent_likely > 0 else 'Less'}",
                    'Purchases per fan vs. gen pop': f"{abs(percent_purch):.1f}% {'More' if percent_purch > 0 else 'Less'}"
                })

        return pd.DataFrame(results)

    def _generate_insights(self, category_stats, subcategory_stats):
        """Generate bullet point insights"""
        insights = []

        # Check if we have valid data
        if category_stats['percent_fans'] == "N/A":
            insights.append(f"No data available for {self.current_category} category")
            return insights

        # Category-level insights
        raw = category_stats['raw_values']

        # Insight 1: Likelihood to spend
        if raw['percent_likely'] > 0:
            insights.append(
                f"Jazz Fans are {raw['percent_likely']:.1f}% MORE likely to spend on "
                f"{self.current_category} than the Utah General Population"
            )
        else:
            insights.append(
                f"Jazz Fans are {abs(raw['percent_likely']):.1f}% LESS likely to spend on "
                f"{self.current_category} than the Utah General Population"
            )

        # Insight 2: Purchase frequency
        if raw['percent_purch'] > 0:
            insights.append(
                f"Jazz Fans make an average of {raw['percent_purch']:.1f}% more purchases on "
                f"{self.current_category} than the Utah General Population"
            )
        else:
            insights.append(
                f"Jazz Fans make an average of {abs(raw['percent_purch']):.1f}% less purchases on "
                f"{self.current_category} than the Utah General Population"
            )

        # Insight 3: Top subcategory
        if not subcategory_stats.empty:
            top_sub = subcategory_stats.iloc[0]
            # Get column names that contain newlines
            likelihood_col = [col for col in subcategory_stats.columns if 'likely' in col][0]
            insights.append(
                f"Jazz Fans are {top_sub[likelihood_col]} "
                f"likely to spend on {top_sub['Subcategory']} than the Utah General Population"
            )

        # Year-over-year analysis
        yoy_insight = self._get_yoy_insights()
        if yoy_insight:
            insights.append(yoy_insight)

        # NBA comparison
        nba_insight = self._get_nba_comparison()
        if nba_insight:
            insights.append(nba_insight)

        # QSR specific insight (if restaurants)
        if self.current_category == "Restaurants":
            qsr_insight = self._get_qsr_insight()
            if qsr_insight:
                insights.append(qsr_insight)

        return insights

    def _get_qsr_insight(self):
        """Get QSR-specific insight for restaurants"""
        query = """
        SELECT * FROM V_UTAH_JAZZ_SIL_SUBCATEGORY_INDEXING_YOY 
        WHERE TRIM(CATEGORY) = %(category)s 
        AND TRIM(SUBCATEGORY) = 'Restaurants - QSR & Fast Casual'
        AND COMPARISON_POPULATION = 'General Population'
        """

        df = query_to_dataframe(query, {'category': self.current_category.strip()})

        if not df.empty:
            df['SPC'] = pd.to_numeric(df['SPC'])
            df['COMPARISON_SPC'] = pd.to_numeric(df['COMPARISON_SPC'])

            # Filter out 2025 data and calculate means
            qsr_stats = df[df['TRANSACTION_YEAR'] != '2025-01-01'].agg({
                'SPC': 'mean',
                'COMPARISON_SPC': 'mean'
            })

            return (f"Jazz Fans spend an average of ${qsr_stats['SPC']:.2f} per fan per year "
                    f"on QSR and Fast Casual Restaurants compared to ${qsr_stats['COMPARISON_SPC']:.2f} "
                    f"per person per year in the Utah General Population")

        return None

    def _get_merchant_stats(self):
        """Get merchant-level statistics"""
        query = f"""
        SELECT * FROM V_UTAH_JAZZ_SIL_MERCHANT_INDEXING_ALL_TIME 
        WHERE TRIM(CATEGORY) = %(category)s
        """

        df = query_to_dataframe(query, {'category': self.current_category.strip()})

        # Check if we have any data
        if df.empty:
            print(f"⚠️  No merchant data found for category: {self.current_category}")
            return pd.DataFrame(), []

        # Get top 5 merchants by audience percentage
        jazz_fans = df[df['AUDIENCE'] == 'Utah Jazz Fans']

        if jazz_fans.empty:
            print(f"⚠️  No Utah Jazz Fans data found in merchant table for category: {self.current_category}")
            return pd.DataFrame(), []

        top_5_merchants = (jazz_fans
                           .sort_values('PERC_AUDIENCE', ascending=False)
                           .drop_duplicates('MERCHANT')
                           .head(5)['MERCHANT'].tolist())

        if not top_5_merchants:
            print(f"⚠️  No merchants found for category: {self.current_category}")
            return pd.DataFrame(), []

        print(f"Found {len(top_5_merchants)} merchants for {self.current_category}")

        # Build merchant comparison table
        merchant_results = []

        for merchant in top_5_merchants:
            jazz_row = df[
                (df['AUDIENCE'] == 'Utah Jazz Fans') &
                (df['MERCHANT'] == merchant) &
                (df['COMPARISON_POPULATION'] == 'Local Gen Pop (Excl. Jazz)')
                ]

            if not jazz_row.empty:
                row = jazz_row.iloc[0]
                percent_fans = float(row['PERC_AUDIENCE']) * 100
                percent_likely = float(row['PERC_INDEX']) - 100
                ppc_diff = self._calculate_percent_diff(
                    float(row['PPC']),
                    float(row['COMPARISON_PPC'])
                )

                merchant_results.append({
                    'Brand': merchant,
                    'Percent of Fans Who Spend': f"{percent_fans:.1f}%",
                    'How likely fans are to spend vs. gen pop':
                        f"{abs(percent_likely):.1f}% {'More' if percent_likely >= 0 else 'Less'}",
                    'Purchases Per Fan (vs. Gen Pop)':
                        f"{abs(ppc_diff):.1f}% {'More' if ppc_diff >= 0 else 'Less'}"
                })

        if not merchant_results:
            print(f"⚠️  No merchant comparison data found")
            return pd.DataFrame(), top_5_merchants

        merchant_df = pd.DataFrame(merchant_results)
        merchant_df['Rank'] = range(1, len(merchant_df) + 1)

        # Reorder columns
        cols = ['Rank', 'Brand', 'Percent of Fans Who Spend',
                'How likely fans are to spend vs. gen pop', 'Purchases Per Fan (vs. Gen Pop)']

        return merchant_df[cols], top_5_merchants

    def _generate_merchant_insights(self, merchant_stats):
        """Generate merchant-specific insights"""
        merchant_df, top_5_merchants = merchant_stats
        insights = []

        # Check if we have any data
        if merchant_df.empty or not top_5_merchants:
            insights.append(f"No merchant data available for {self.current_category} category")
            return insights

        # Top merchant insight
        top_merchant = merchant_df.iloc[0]
        insights.append(
            f"{top_merchant['Percent of Fans Who Spend']} of Utah Jazz Fans spent at {top_merchant['Brand']}"
        )

        # Get year-over-year merchant data for additional insights
        yoy_query = f"""
        SELECT * FROM V_UTAH_JAZZ_SIL_MERCHANT_INDEXING_YOY 
        WHERE TRIM(CATEGORY) = %(category)s AND AUDIENCE = 'Utah Jazz Fans'
        """

        yoy_df = query_to_dataframe(yoy_query, {'category': self.current_category.strip()})

        if not yoy_df.empty:
            yoy_df['PPC'] = pd.to_numeric(yoy_df['PPC'])
            yoy_df['SPC'] = pd.to_numeric(yoy_df['SPC'])

            # Average purchases per year
            ppc_summary = (yoy_df[yoy_df['MERCHANT'].isin(top_5_merchants)]
                           .groupby('MERCHANT')['PPC']
                           .mean()
                           .sort_values(ascending=False))

            if not ppc_summary.empty:
                top_ppc_merchant = ppc_summary.index[0]
                avg_ppc = ppc_summary.iloc[0]
                insights.append(
                    f"Utah Jazz Fans average {avg_ppc:.1f} purchases per year at {top_ppc_merchant}"
                )

            # Average spend per year
            spc_summary = (yoy_df[yoy_df['MERCHANT'].isin(top_5_merchants)]
                           .groupby('MERCHANT')['SPC']
                           .mean()
                           .sort_values(ascending=False))

            if not spc_summary.empty:
                top_spc_merchant = spc_summary.index[0]
                avg_spc = spc_summary.iloc[0]
                insights.append(
                    f"Utah Jazz Fans spent an average of ${avg_spc:.2f} per year at {top_spc_merchant}"
                )

        # NBA comparison for merchants
        nba_comparison = self._get_merchant_nba_comparison(top_5_merchants)
        if nba_comparison:
            insights.append(nba_comparison)

        # Composite index recommendation
        composite_target = self._get_composite_recommendation(top_5_merchants)
        if composite_target:
            insights.append(
                f"The Utah Jazz should target {composite_target['merchant']} as a sponsor "
                f"based on a composite index of {composite_target['index']:.0f}"
            )

        return insights

    def create_powerpoint(self, analysis_results, output_file='utah_jazz_sponsor_analysis.pptx'):
        """Create PowerPoint presentation with analysis results"""
        # Load template or create new presentation
        template_path = Path('SIL_Template.pptx')
        if template_path.exists():
            prs = Presentation(str(template_path))
            print("✅ Loaded SIL template")
        else:
            print("⚠️  Template not found, creating new presentation")
            prs = Presentation()

        # Set slide size to 16:9 widescreen if creating new presentation
        if not template_path.exists():
            prs.slide_width = Inches(13.333)  # 16:9 aspect ratio width
            prs.slide_height = Inches(7.5)  # 16:9 aspect ratio height

        # Add category analysis slide
        self._add_category_slide(prs, analysis_results)

        # Add merchant analysis slide
        self._add_merchant_slide(prs, analysis_results)

        # Save presentation
        prs.save(output_file)
        print(f"✅ PowerPoint saved as: {output_file}")

        return output_file

    def _add_category_slide(self, prs, results):
        """Add category analysis slide with proper formatting"""
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)

        # Clear existing placeholders
        for shape in slide.placeholders:
            sp = shape.element
            sp.getparent().remove(sp)

        # Add title with Red Hat Display font
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(7), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = f"{self.current_category} Sponsor Analysis"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(24)
        title_para.font.bold = True
        # Apply font to both paragraph and run level
        title_para.font.name = 'Red Hat Display'
        for run in title_para.runs:
            run.font.name = 'Red Hat Display'
            run.font.bold = True

        # Add headers
        self._add_slide_headers(slide)

        # Create category stats table
        cat_stats = results['category_stats']
        cat_data = [
            ['Category', self.current_category],
            ['Percent of Fans Who Spend', cat_stats['percent_fans']],
            ['How likely fans are to spend vs. gen pop', cat_stats['likelihood']],
            ['Purchases per fan vs. gen pop', cat_stats['purchases']]
        ]

        # Add category table
        self._add_table_to_slide(
            slide, cat_data,
            left=Inches(3.7), top=Inches(1.2),
            width=Inches(5.2), height=Inches(1.0)
        )

        # Add subcategory table
        sub_df = results['subcategory_stats']
        if not sub_df.empty:
            # Convert DataFrame to list format for table
            sub_data = [sub_df.columns.tolist()] + sub_df.values.tolist()

            self._add_table_to_slide(
                slide, sub_data,
                left=Inches(3.7), top=Inches(2.5),
                width=Inches(5.2), height=Inches(2.5)
            )

        # Add insights as bullets with Red Hat Display font
        insights_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.3), Inches(2.8), Inches(4.5)
        )
        tf = insights_box.text_frame
        tf.word_wrap = True

        for i, insight in enumerate(results['insights']):
            if i > 0:
                p = tf.add_paragraph()
                p.text = ""  # Empty line between bullets

            p = tf.add_paragraph()
            p.text = f"• {insight}"
            p.font.size = Pt(10)
            p.font.name = 'Red Hat Display'
            # Apply to runs as well
            for run in p.runs:
                run.font.name = 'Red Hat Display'

    def _add_merchant_slide(self, prs, results):
        """Add merchant analysis slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # Clear placeholders
        for shape in slide.placeholders:
            sp = shape.element
            sp.getparent().remove(sp)

        # Title with Red Hat Display font
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(7), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = f"{self.current_category} Sponsor Analysis"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(24)
        title_para.font.bold = True
        self._apply_font_to_text(title_para, 'bold')

        # Add headers
        self._add_slide_headers(slide)

        # Add merchant table
        merchant_df, _ = results['merchant_stats']
        if not merchant_df.empty:
            # The merchant_df already has clean column names from _get_merchant_stats
            # Just convert to list format
            merchant_data = [merchant_df.columns.tolist()] + merchant_df.values.tolist()

            self._add_table_to_slide(
                slide, merchant_data,
                left=Inches(3.7), top=Inches(1.2),
                width=Inches(5.2), height=Inches(2.5)
            )

        # Add merchant insights with Red Hat Display font
        insights_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.3), Inches(2.8), Inches(5.0)
        )
        tf = insights_box.text_frame
        tf.word_wrap = True

        merchant_insights = results['merchant_insights']

        # Add first 4 bullets
        for i, insight in enumerate(merchant_insights[:4]):
            if i > 0:
                p = tf.add_paragraph()
                p.text = ""  # Empty line

            p = tf.add_paragraph()
            p.text = f"• {insight}"
            p.font.size = Pt(10)
            self._apply_font_to_text(p)

        # Add Top Brand Target section if exists
        if len(merchant_insights) > 4:
            # Add spacing and header
            p = tf.add_paragraph()
            p.text = ""

            p = tf.add_paragraph()
            p.text = "Top Brand Target"
            p.font.bold = True
            p.font.size = Pt(11)
            self._apply_font_to_text(p, 'bold')

            p = tf.add_paragraph()
            p.text = ""

            # Add recommendation
            p = tf.add_paragraph()
            p.text = f"• {merchant_insights[4]}"
            p.font.size = Pt(10)
            self._apply_font_to_text(p)

            # Add explanation
            p = tf.add_paragraph()
            p.text = ""

            p = tf.add_paragraph()
            p.text = ("The composite index is a Sports Innovation Lab score that takes into "
                      "account multiple fan spending metrics including how many fans spend, "
                      "how often they spend, and how much they spend.")
            p.font.size = Pt(8)
            p.font.italic = True
            self._apply_font_to_text(p)

    def _add_table_to_slide(self, slide, data, left, top, width, height):
        """Add a formatted table to the slide with Red Hat Display font"""
        rows = len(data)
        cols = len(data[0]) if data else 0

        if rows == 0 or cols == 0:
            return

        # Add table
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        # Format table
        for i, row_data in enumerate(data):
            for j, cell_value in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = str(cell_value)

                # Get the paragraph
                paragraph = cell.text_frame.paragraphs[0]

                # Format header row
                if i == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(242, 242, 242)
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(10)
                    paragraph.alignment = PP_ALIGN.CENTER
                    self._apply_font_to_text(paragraph, 'bold')
                else:
                    paragraph.font.size = Pt(10)
                    # Right-align numeric columns
                    if j > 1:
                        paragraph.alignment = PP_ALIGN.RIGHT
                    self._apply_font_to_text(paragraph)

        # Set column widths (autofit equivalent)
        for i, col in enumerate(table.columns):
            if i == 0:  # Rank column
                col.width = Inches(0.5)
            elif i == 1:  # Brand column
                col.width = Inches(1.5)
            else:  # Other columns
                col.width = Inches(1.0)

    def _add_slide_headers(self, slide):
        """Add header/footer text to slides with Red Hat Display font"""
        # Left header
        left_box = slide.shapes.add_textbox(
            Inches(0.2), Inches(0.1), Inches(3), Inches(0.3)
        )
        left_frame = left_box.text_frame
        left_frame.text = "Utah Jazz"
        p = left_frame.paragraphs[0]
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(128, 128, 128)
        p.font.italic = True
        self._apply_font_to_text(p)

        # Right header
        right_box = slide.shapes.add_textbox(
            Inches(7), Inches(0.1), Inches(3), Inches(0.3)
        )
        right_frame = right_box.text_frame
        right_frame.text = f"Sponsor Spending Analysis: {self.current_category}"
        p = right_frame.paragraphs[0]
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(128, 128, 128)
        p.font.italic = True
        p.alignment = PP_ALIGN.RIGHT
        self._apply_font_to_text(p)

    # Helper methods
    def _calculate_percent_diff(self, value1, value2):
        """Calculate percentage difference"""
        if value2 == 0:
            return 0
        return ((value1 - value2) / value2) * 100

    def _get_yoy_insights(self):
        """Get year-over-year insights"""
        query = f"""
        SELECT TRANSACTION_YEAR, PERC_AUDIENCE 
        FROM V_UTAH_JAZZ_SIL_CATEGORY_INDEXING_YOY 
        WHERE TRIM(CATEGORY) = %(category)s
        AND TRANSACTION_YEAR IN ('2023-01-01', '2024-01-01')
        """

        df = query_to_dataframe(query, {'category': self.current_category.strip()})

        if df.empty or len(df) < 2:
            return None

        df['PERC_AUDIENCE'] = pd.to_numeric(df['PERC_AUDIENCE'])
        df = df.sort_values('TRANSACTION_YEAR')

        pct_change = ((df.iloc[1]['PERC_AUDIENCE'] - df.iloc[0]['PERC_AUDIENCE']) /
                      df.iloc[0]['PERC_AUDIENCE'] * 100)

        if pct_change > 0:
            return (f"{self.current_category} saw an INCREASE of {pct_change:.1f}% "
                    f"of Jazz fans spending on the category from 2023-2024")

        return None

    def _get_nba_comparison(self):
        """Get NBA fan comparison"""
        query = f"""
        SELECT * FROM V_UTAH_JAZZ_SIL_SUBCATEGORY_INDEXING_ALL_TIME 
        WHERE TRIM(CATEGORY) = %(category)s
        AND COMPARISON_POPULATION = 'NBA Fans'
        ORDER BY PERC_INDEX ASC
        LIMIT 1
        """

        df = query_to_dataframe(query, {'category': self.current_category.strip()})

        if not df.empty:
            row = df.iloc[0]
            likely = float(row['PERC_INDEX']) - 100

            if likely > 0:
                return (f"Jazz Fans are {likely:.1f}% MORE likely to spend on "
                        f"{row['SUBCATEGORY']} than NBA Fans")
            else:
                return (f"Jazz Fans are {abs(likely):.1f}% LESS likely to spend on "
                        f"{row['SUBCATEGORY']} than NBA Fans")

        return None

    def _get_merchant_nba_comparison(self, top_merchants):
        """Get merchant NBA comparison"""
        # Fixed: Build merchant list for SQL IN clause - properly escape single quotes
        merchant_list = ', '.join(["'{}'".format(m.replace("'", "''")) for m in top_merchants])

        query = f"""
        SELECT * FROM V_UTAH_JAZZ_SIL_MERCHANT_INDEXING_ALL_TIME 
        WHERE COMPARISON_POPULATION = 'NBA Fans' 
        AND MERCHANT IN ({merchant_list})
        """

        df = query_to_dataframe(query)

        if df.empty:
            return None

        # Find highest index
        index_cols = ['PERC_INDEX', 'SPC_INDEX', 'SPP_INDEX', 'PPC_INDEX']
        for col in index_cols:
            df[col] = pd.to_numeric(df[col])

        # Find max value and its location
        max_val = df[index_cols].max().max()
        max_loc = df[index_cols].stack().idxmax()

        merchant = df.loc[max_loc[0], 'MERCHANT']
        metric = max_loc[1]

        if metric == 'PERC_INDEX':
            percent_more = max_val - 100
            return f"Utah Jazz Fans are {percent_more:.1f}% more likely to spend on {merchant} than NBA Fans"
        elif metric == 'PPC_INDEX':
            row = df.loc[max_loc[0]]
            ppc_diff = self._calculate_percent_diff(
                float(row['PPC']),
                float(row['COMPARISON_PPC'])
            )
            return f"Utah Jazz Fans make an average of {ppc_diff:.1f}% more purchases on {merchant} than NBA Fans"

        return None

    def _get_composite_recommendation(self, top_merchants):
        """Get composite index recommendation"""
        # Fixed: Build merchant list for SQL IN clause - properly escape single quotes
        merchant_list = ', '.join(["'{}'".format(m.replace("'", "''")) for m in top_merchants])

        query = f"""
        SELECT MERCHANT, COMPOSITE_INDEX 
        FROM V_UTAH_JAZZ_SIL_MERCHANT_INDEXING_ALL_TIME 
        WHERE AUDIENCE = 'Utah Jazz Fans' 
        AND COMPARISON_POPULATION = 'Local Gen Pop (Excl. Jazz)'
        AND MERCHANT IN ({merchant_list})
        ORDER BY COMPOSITE_INDEX DESC
        LIMIT 1
        """

        df = query_to_dataframe(query)

        if not df.empty:
            return {
                'merchant': df.iloc[0]['MERCHANT'],
                'index': float(df.iloc[0]['COMPOSITE_INDEX'])
            }

        return None


def main():
    """Main execution function"""
    # Initialize analyzer
    analyzer = SponsorshipAnalyzer()

    try:
        # Test connection first
        print("🔄 Testing Snowflake connection...")
        test_conn = get_connection()
        print("✅ Connection successful!")

        # Check available data
        available_categories = analyzer.check_available_data()

        # Analyze specific category (matching R script)
        category = "Restaurants"
        results = analyzer.analyze_category(category)

        # Display results
        print("\n📊 Category Statistics:")
        print(f"  - Percent of fans who spend: {results['category_stats']['percent_fans']}")
        print(f"  - Likelihood vs gen pop: {results['category_stats']['likelihood']}")
        print(f"  - Purchases vs gen pop: {results['category_stats']['purchases']}")

        print("\n📊 Top Subcategories:")
        if not results['subcategory_stats'].empty:
            print(results['subcategory_stats'].to_string(index=False))
        else:
            print("  No subcategory data available")

        print("\n💡 Key Insights:")
        for i, insight in enumerate(results['insights'], 1):
            print(f"  {i}. {insight}")

        print("\n🏪 Top Merchants:")
        merchant_df, _ = results['merchant_stats']
        if not merchant_df.empty:
            print(merchant_df.to_string(index=False))
        else:
            print("  No merchant data available")

        print("\n💡 Merchant Insights:")
        for i, insight in enumerate(results['merchant_insights'], 1):
            print(f"  {i}. {insight}")

        # Create PowerPoint
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_file = f"utah_jazz_{category.lower()}_analysis_{timestamp}.pptx"
        analyzer.create_powerpoint(results, output_file)

        print(f"\n✅ Analysis complete! PowerPoint saved as: {output_file}")

        # Process all categories if desired
        process_all = input("\nProcess all categories? (y/n): ")
        if process_all.lower() == 'y':
            for cat in analyzer.categories:
                if cat != category:  # Skip already processed
                    print(f"\n{'=' * 60}")
                    try:
                        cat_results = analyzer.analyze_category(cat)
                        cat_output = f"utah_jazz_{cat.lower()}_analysis_{timestamp}.pptx"
                        analyzer.create_powerpoint(cat_results, cat_output)
                        print(f"✅ {cat} analysis saved as: {cat_output}")
                    except Exception as e:
                        print(f"❌ Error processing {cat}: {str(e)}")

    except Exception as e:
        print(f"\n❌ Error: {str(e)}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    main()