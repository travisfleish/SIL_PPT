"""
PowerPoint Automation Framework for Sports Innovation Lab
Generates branded insight reports from fan data
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
import pandas as pd
import os
from datetime import datetime

class SportsInsightsGenerator:
    def __init__(self, template_path=None):
        """Initialize with optional template."""
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()

        # Brand colors
        self.JAZZ_BLUE = RGBColor(29, 66, 138)  # #1D428A
        self.JAZZ_YELLOW = RGBColor(255, 199, 44)  # #FFC72C
        self.JAZZ_GREEN = RGBColor(0, 43, 92)  # #002B5C

    def add_title_slide(self, team_name, report_date=None):
        """Add branded title slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])

        title = slide.shapes.title
        title.text = f"{team_name}\nSponsorship Insights Report"

        # Format title
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(44)
                run.font.bold = True
                run.font.color.rgb = self.JAZZ_BLUE

        # Add date
        if report_date:
            subtitle = slide.placeholders[1]
            subtitle.text = report_date.strftime("%B %Y")

    def add_fan_wheel_slide(self, team_name, wheel_image_path,
                           audience_data, fan_description):
        """Add slide with fan wheel and audience index chart."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Blank

        # Add header
        header_left = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(4), Inches(0.5)
        )
        header_left.text = team_name
        header_left.text_frame.paragraphs[0].font.size = Pt(18)
        header_left.text_frame.paragraphs[0].font.bold = True

        header_right = slide.shapes.add_textbox(
            Inches(5.5), Inches(0.3), Inches(4), Inches(0.5)
        )
        header_right.text = f"Fan Behaviors: How Are {team_name} Fans Unique"
        header_right.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        header_right.text_frame.paragraphs[0].font.size = Pt(18)

        # Add wheel image
        if os.path.exists(wheel_image_path):
            slide.shapes.add_picture(
                wheel_image_path,
                Inches(0.3), Inches(1.2),
                width=Inches(4.5)
            )

        # Add audience index chart
        self._add_audience_index_chart(slide, audience_data)

        # Add description
        desc_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(5.8), Inches(9.4), Inches(0.8)
        )
        desc_box.text = fan_description
        desc_box.text_frame.paragraphs[0].font.size = Pt(16)
        desc_box.text_frame.paragraphs[0].font.bold = True

    def _add_audience_index_chart(self, slide, audience_data):
        """Add horizontal bar chart for audience index."""
        # Create chart data
        chart_data = CategoryChartData()

        # Sort data by value descending
        sorted_data = sorted(audience_data.items(),
                           key=lambda x: x[1], reverse=True)

        categories = [item[0] for item in sorted_data]
        values = [item[1] for item in sorted_data]

        chart_data.categories = categories
        chart_data.add_series('% Audience Index', values)

        # Add chart to slide
        x, y = Inches(5), Inches(1.5)
        cx, cy = Inches(4.5), Inches(3.5)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
        ).chart

        # Format chart
        chart.has_legend = False
        chart.has_title = True
        chart.chart_title.text_frame.text = "% Audience Index"

        # Format value axis to show percentages
        value_axis = chart.value_axis
        value_axis.maximum_scale = 100
        value_axis.has_major_gridlines = True

        # Color bars based on value
        series = chart.series[0]
        for idx, point in enumerate(series.points):
            if values[idx] >= 60:
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = self.JAZZ_YELLOW
            else:
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = RGBColor(128, 128, 128)

        # Add value labels
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(10)
        data_labels.font.bold = True
        data_labels.position = XL_DATA_LABEL_POSITION.INSIDE_END

    def add_category_slide(self, category_name, brand_data, insights):
        """Add slide for specific category analysis."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
        )
        title_box.text = f"{category_name} Sponsor Analysis"
        title_box.text_frame.paragraphs[0].font.size = Pt(24)
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.color.rgb = self.JAZZ_BLUE

        # Add data table or visualization
        self._add_brand_comparison_chart(slide, brand_data)

        # Add insights
        insights_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5), Inches(9), Inches(2)
        )
        insights_box.text = insights
        insights_box.text_frame.word_wrap = True

        for paragraph in insights_box.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.space_after = Pt(12)

    def _add_brand_comparison_chart(self, slide, brand_data):
        """Add chart comparing brands in category."""
        chart_data = CategoryChartData()

        brands = list(brand_data.keys())
        chart_data.categories = brands

        # Multiple series for different metrics
        metrics = ['% Fanbase', 'Avg Annual Spend', 'Transactions']
        for metric in metrics:
            values = [brand_data[brand].get(metric, 0) for brand in brands]
            chart_data.add_series(metric, values)

        x, y = Inches(0.5), Inches(1.5)
        cx, cy = Inches(9), Inches(3)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    def save(self, filename):
        """Save the presentation."""
        self.prs.save(filename)
        print(f"✅ Presentation saved as {filename}")

# Example usage
if __name__ == "__main__":
    # Initialize generator
    generator = SportsInsightsGenerator()

    # Add title slide
    generator.add_title_slide("Utah Jazz", datetime.now())

    # Add fan wheel slide
    audience_data = {
        'Live Entertainment Seekers': 71,
        'Cost Conscious': 70,
        'Travelers': 62,
        'Gen Z Brand Shoppers': 60,
        'Beauty Enthusiasts': 45,
        'Movie Buffs': 45,
        'Fans of Men\'s Sports': 45,
        'Sports Streamer': 39,
        'Gamers': 35,
        'Pet Owners': 32
    }

    generator.add_fan_wheel_slide(
        "Utah Jazz",
        "professional_fan_wheel.png",
        audience_data,
        "Jazz fans are values-driven live entertainment seekers who are on the lookout for a deal and a good movie!"
    )

    # Add category slides
    auto_brands = {
        'AutoZone': {'% Fanbase': 15, 'Avg Annual Spend': 450, 'Transactions': 3.2},
        'Jiffy Lube': {'% Fanbase': 12, 'Avg Annual Spend': 380, 'Transactions': 2.8},
        'Midas': {'% Fanbase': 8, 'Avg Annual Spend': 420, 'Transactions': 2.1}
    }

    generator.add_category_slide(
        "Auto Services",
        auto_brands,
        "AutoZone shows the highest penetration among Jazz fans at 15%, with strong average annual spend. This represents a prime sponsorship opportunity given the high transaction frequency."
    )

    # Save presentation
    generator.save("Utah_Jazz_Insights_Report.pptx")