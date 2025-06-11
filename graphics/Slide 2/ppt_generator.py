#!/usr/bin/env python3
"""
Utah Jazz Fan Insights Slide Generator
Orchestrates the creation of a complete PowerPoint slide with fan wheel and community chart
Fixed version with correct path handling
"""

import os
import sys
import subprocess
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime
import time
from pathlib import Path
from openai import OpenAI
from dotenv import load_dotenv
import pandas as pd
import glob
import json
import re


class JazzInsightsSlideGenerator:
    """Generate a complete Utah Jazz insights slide matching the provided design"""

    def __init__(self):
        # Brand colors
        self.JAZZ_BLUE = RGBColor(29, 66, 138)  # #1D428A
        self.JAZZ_YELLOW = RGBColor(255, 199, 44)  # #FFC72C
        self.JAZZ_GREEN = RGBColor(0, 43, 92)  # #002B5C

        # Load environment variables
        load_dotenv()

        # Initialize OpenAI client for summary generation
        self.openai_client = None
        if os.getenv("OPENAI_API_KEY"):
            self.openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

        # Determine the base directory
        self.current_file = Path(__file__).resolve()
        self.current_dir = self.current_file.parent

        # Check if we're already in 'graphics/Slide 2' directory
        if self.current_dir.name == 'Slide 2' and self.current_dir.parent.name == 'graphics':
            # We're in the graphics/Slide 2 directory
            self.base_dir = self.current_dir.parent.parent  # Go up to project root
            self.scripts_dir = self.current_dir
        else:
            # We're somewhere else (maybe project root)
            self.base_dir = self.current_dir
            self.scripts_dir = self.base_dir / 'graphics' / 'Slide 2'

        print(f"📁 Working directories:")
        print(f"   Base directory: {self.base_dir}")
        print(f"   Scripts directory: {self.scripts_dir}")
        print(f"   Current directory: {os.getcwd()}")

        # File paths
        self.wheel_image_path = None
        self.chart_image_path = None

    def run_community_chart(self):
        """Execute community_per_chart.py to generate the audience index chart"""
        print("\n📊 Step 1: Generating Community Audience Index Chart...")
        print("-" * 50)

        # Find the script
        script_path = self.scripts_dir / "community_per_chart.py"

        if not script_path.exists():
            print(f"❌ Script not found at: {script_path}")
            print(f"   Looking for alternative locations...")

            # Try alternative paths
            alt_paths = [
                self.current_dir / "community_per_chart.py",
                self.base_dir / "community_per_chart.py",
                Path("community_per_chart.py")
            ]

            for alt_path in alt_paths:
                if alt_path.exists():
                    script_path = alt_path
                    print(f"   ✓ Found at: {script_path}")
                    break
            else:
                print("   ❌ Could not find community_per_chart.py")
                return False

        try:
            print(f"   Running: {script_path}")

            # Run the community chart script
            result = subprocess.run(
                [sys.executable, str(script_path)],
                capture_output=True,
                text=True,
                cwd=str(self.base_dir)  # Run from base directory
            )

            if result.returncode == 0:
                print("✅ Community chart generated successfully!")

                # Parse the output to save community data
                self._save_community_data_from_output(result.stdout)

                # Find the generated chart file (latest PNG with timestamp)
                import glob

                # Look in base directory for generated files
                os.chdir(self.base_dir)
                chart_files = glob.glob("jazz_audience_index_*.png")

                if chart_files:
                    # Get the most recent file
                    self.chart_image_path = max(chart_files, key=os.path.getctime)
                    print(f"   Chart saved as: {self.chart_image_path}")
                else:
                    # Fallback to default name
                    if os.path.exists("audience_index_chart.png"):
                        self.chart_image_path = "audience_index_chart.png"
                        print(f"   Chart saved as: {self.chart_image_path}")
                    else:
                        print("   ⚠️  Warning: Chart file not found after generation")
                        return False
                return True
            else:
                print(f"❌ Error generating community chart:")
                print(f"STDOUT: {result.stdout}")
                print(f"STDERR: {result.stderr}")
                return False

        except Exception as e:
            print(f"❌ Failed to run community chart script: {str(e)}")
            return False

    def _save_community_data_from_output(self, output):
        """Parse the script output to extract community data"""
        try:
            community_data = []
            lines = output.split('\n')

            # Look for the data summary section
            in_data_section = False
            for line in lines:
                if "Top communities by composite index" in line:
                    in_data_section = True
                    continue

                if in_data_section and "•" in line:
                    # Parse lines like: "• Live Entertainment Seekers: Composite=2500, Index=71%, Audience=21.0%"
                    parts = line.split(":")
                    if len(parts) >= 2:
                        name = parts[0].replace("•", "").strip()

                        # Extract percentages
                        index_match = re.search(r'Index=(\d+)%', line)
                        audience_match = re.search(r'Audience=([\d.]+)%', line)

                        if index_match:
                            community_data.append({
                                "name": name,
                                "index": int(index_match.group(1)),
                                "audience": float(audience_match.group(1)) if audience_match else 0
                            })

            # Save to JSON for the AI summary generator
            if community_data:
                with open("community_data.json", "w") as f:
                    json.dump(community_data, f, indent=2)
                print(f"   💾 Saved {len(community_data)} communities to community_data.json")

        except Exception as e:
            print(f"   ⚠️  Could not parse community data: {e}")

    def run_fan_wheel(self):
        """Execute generate_fan_wheel_with_logos.py to generate the fan wheel"""
        print("\n🎨 Step 2: Generating Fan Wheel Visualization...")
        print("-" * 50)

        # Find the script
        script_path = self.scripts_dir / "generate_fan_wheel_with_logos.py"

        if not script_path.exists():
            print(f"❌ Script not found at: {script_path}")
            print(f"   Looking for alternative locations...")

            # Try alternative paths
            alt_paths = [
                self.current_dir / "generate_fan_wheel_with_logos.py",
                self.base_dir / "generate_fan_wheel_with_logos.py",
                Path("generate_fan_wheel_with_logos.py")
            ]

            for alt_path in alt_paths:
                if alt_path.exists():
                    script_path = alt_path
                    print(f"   ✓ Found at: {script_path}")
                    break
            else:
                print("   ❌ Could not find generate_fan_wheel_with_logos.py")
                return False

        try:
            print(f"   Running: {script_path}")

            # Run the fan wheel script
            result = subprocess.run(
                [sys.executable, str(script_path)],
                capture_output=True,
                text=True,
                cwd=str(self.base_dir)  # Run from base directory
            )

            if result.returncode == 0:
                print("✅ Fan wheel generated successfully!")

                # Check for the generated file
                wheel_path = self.base_dir / "professional_fan_wheel.png"
                if wheel_path.exists():
                    self.wheel_image_path = str(wheel_path)
                    print(f"   Wheel saved as: {self.wheel_image_path}")
                else:
                    print("   ⚠️  Warning: Fan wheel file not found after generation")
                    return False
                return True
            else:
                print(f"❌ Error generating fan wheel:")
                print(f"STDOUT: {result.stdout}")
                print(f"STDERR: {result.stderr}")
                return False

        except Exception as e:
            print(f"❌ Failed to run fan wheel script: {str(e)}")
            return False

    def generate_ai_summary(self):
        """Generate AI-powered summary based on available data"""
        if not self.openai_client:
            print("⚠️  OpenAI API key not found, using default summary")
            return "Jazz fans are values-driven live entertainment seekers who are on the lookout for a deal and a good movie!"

        print("\n🤖 Generating AI-powered fan summary...")

        try:
            # Get actual community data from the chart
            community_data = []
            chart_files = glob.glob("jazz_audience_index_*.png")

            # Read the actual Snowflake data if we can
            actual_communities = []
            try:
                # Try to read from a saved results file or query result
                import json
                if os.path.exists("community_data.json"):
                    with open("community_data.json", "r") as f:
                        community_data = json.load(f)
                        actual_communities = [f"{item['name']} ({item['index']}%)" for item in community_data[:5]]
            except:
                # Use what we can see from the chart
                actual_communities = [
                    "NHL fans (21%)",
                    "Hockey fans (23%)",
                    "Skiers (2%)",
                    "Live Entertainment Seekers (25%)",
                    "Surf fans (1%)"
                ]

            # Get fan wheel data
            wheel_brands = []
            wheel_behaviors = []
            if os.path.exists("mock_fan_wheel.csv"):
                df = pd.read_csv("mock_fan_wheel.csv")
                wheel_brands = df['brand'].tolist()
                wheel_behaviors = df['behavior'].str.replace('\n', ' ').tolist()

            # Create a detailed prompt
            prompt = f"""Analyze this Utah Jazz fan data and write ONE engaging sentence that accurately reflects what the data shows.

Community Index Data (from chart):
{chr(10).join(actual_communities)}

Brand Affinities (from wheel):
{chr(10).join([f"- {behavior}" for behavior in wheel_behaviors[:5]])}

Rules:
1. The sentence MUST reference specific insights from the data
2. Focus on 2-3 KEY traits that are actually shown in the data
3. Don't make up traits not supported by the data
4. Keep it under 25 words
5. End with an exclamation point
6. Be specific - mention actual communities or behaviors shown

Example: "Jazz fans are [specific trait from data] who [specific behavior from wheel] and [another trait from chart]!"
"""

            response = self.openai_client.chat.completions.create(
                model="gpt-4",
                messages=[
                    {"role": "system",
                     "content": "You are a data analyst who creates accurate summaries based only on the data provided. Never invent traits not shown in the data."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.3,  # Lower temperature for more accurate/factual
                max_tokens=60
            )

            summary = response.choices[0].message.content.strip()
            print(f"✅ Generated: {summary}")

            # Save the summary for reference
            with open("generated_summary.txt", "w") as f:
                f.write(summary)

            return summary

        except Exception as e:
            print(f"⚠️  Error generating AI summary: {e}")
            return "Jazz fans are values-driven live entertainment seekers who are on the lookout for a deal and a good movie!"

    def create_powerpoint_slide(self, output_filename="Utah_Jazz_Fan_Insights.pptx", custom_summary=None):
        """Create the PowerPoint slide with the generated visualizations"""
        print("\n📝 Step 3: Creating PowerPoint Slide...")
        print("-" * 50)

        # Ensure we're working from base directory
        os.chdir(self.base_dir)

        # Check if images exist
        if not self.wheel_image_path or not os.path.exists(self.wheel_image_path):
            print(f"❌ Fan wheel image not found: {self.wheel_image_path}")
            return False

        if not self.chart_image_path or not os.path.exists(self.chart_image_path):
            print(f"❌ Community chart image not found: {self.chart_image_path}")
            return False

        try:
            # Create presentation
            prs = Presentation()

            # Try to use the completely blank layout (usually index 6)
            # If not available, use blank layout and remove placeholders
            try:
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # Completely blank
            except:
                # Fallback to blank layout
                slide = prs.slides.add_slide(prs.slide_layouts[5])

                # Remove ALL placeholders including title placeholders
                shapes_to_remove = []
                for shape in slide.shapes:
                    if shape.is_placeholder:
                        shapes_to_remove.append(shape)

                # Remove the placeholders
                for shape in shapes_to_remove:
                    sp = shape.element
                    sp.getparent().remove(sp)

            # 1. Add header - Left side (Utah Jazz)
            header_left = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.2), Inches(2), Inches(0.5)
            )
            tf = header_left.text_frame
            tf.text = "Utah Jazz"
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.name = "Arial"

            # 2. Add header - Right side
            header_right = slide.shapes.add_textbox(
                Inches(4.5), Inches(0.2), Inches(5), Inches(0.5)
            )
            tf = header_right.text_frame
            tf.text = "Fan Behaviors: How Are Utah Jazz Fans Unique"
            tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.name = "Arial"

            # 3. Add fan wheel image (left side)
            # Calculate vertical centering
            # Slide height is 7.5", headers take ~0.7", description takes ~1.0" at bottom
            # Available space is roughly from 0.7" to 6.5" = 5.8" of space
            # We need to estimate image heights to center them

            # Assuming both images are roughly square when displayed
            wheel_left = Inches(0.3)
            wheel_width = Inches(4.2)
            # Estimate wheel height (assuming square aspect ratio)
            wheel_height_estimate = wheel_width

            # Center vertically: (7.5 - wheel_height) / 2
            wheel_top = Inches((7.5 - 4.2) / 2)  # This centers the wheel vertically on the slide

            slide.shapes.add_picture(
                self.wheel_image_path,
                wheel_left, wheel_top,
                width=wheel_width
            )
            print(f"   ✓ Added fan wheel at position ({wheel_left}, {wheel_top}) - vertically centered")

            # 4. Add community chart image (right side)
            chart_left = Inches(4.8)
            chart_width = Inches(4.8)
            # For the chart, we need to match its vertical center with the wheel's center
            # Assuming the chart is roughly 3.5" tall based on typical aspect ratio
            chart_height_estimate = Inches(3.5)

            # Center vertically to match wheel's center
            chart_top = Inches((7.5 - 3.5) / 2)

            slide.shapes.add_picture(
                self.chart_image_path,
                chart_left, chart_top,
                width=chart_width
            )
            print(f"   ✓ Added community chart at position ({chart_left}, {chart_top}) - vertically centered")

            # 5. Add description text at bottom
            # Generate AI summary if not provided
            if custom_summary:
                description_text = custom_summary
            else:
                description_text = self.generate_ai_summary()

            desc_box = slide.shapes.add_textbox(
                Inches(0.3), Inches(6.3), Inches(9.4), Inches(0.8)  # Near bottom of slide
            )
            tf = desc_box.text_frame
            tf.text = description_text
            tf.word_wrap = True

            # Format the description text
            para = tf.paragraphs[0]
            para.font.size = Pt(16)
            para.font.bold = True
            para.font.name = "Arial"
            para.alignment = PP_ALIGN.LEFT

            # Save presentation
            output_path = self.base_dir / output_filename
            prs.save(str(output_path))
            print(f"\n✅ PowerPoint presentation saved as: {output_path}")

            # Create a summary
            print("\n📋 Slide Summary:")
            print(f"   - Slide dimensions: 10\" x 7.5\" (standard)")
            print(f"   - Fan wheel: {wheel_width} wide")
            print(f"   - Community chart: {chart_width} wide")
            print(f"   - Total shapes added: {len(slide.shapes)}")

            return True

        except Exception as e:
            print(f"❌ Failed to create PowerPoint slide: {str(e)}")
            import traceback
            traceback.print_exc()
            return False

    def generate_complete_slide(self, output_filename="Utah_Jazz_Fan_Insights.pptx", use_ai_summary=True):
        """Main orchestration function"""
        print("🏀 Utah Jazz Fan Insights Slide Generator")
        print("=" * 60)
        print(f"Started at: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

        start_time = time.time()

        # Step 1: Generate community chart
        if not self.run_community_chart():
            print("\n❌ Failed to generate community chart. Exiting.")
            return False

        # Step 2: Generate fan wheel
        if not self.run_fan_wheel():
            print("\n❌ Failed to generate fan wheel. Exiting.")
            return False

        # Step 3: Create PowerPoint slide
        custom_summary = None
        if use_ai_summary and self.openai_client:
            # Generate AI summary before creating slide
            custom_summary = self.generate_ai_summary()

        if not self.create_powerpoint_slide(output_filename, custom_summary):
            print("\n❌ Failed to create PowerPoint slide. Exiting.")
            return False

        # Calculate total time
        total_time = time.time() - start_time

        print("\n" + "=" * 60)
        print("✅ PROCESS COMPLETE!")
        print(f"Total time: {total_time:.2f} seconds")
        print(f"\nGenerated files in {self.base_dir}:")
        print(f"  1. Fan wheel: {os.path.basename(self.wheel_image_path)}")
        print(f"  2. Community chart: {os.path.basename(self.chart_image_path)}")
        print(f"  3. PowerPoint: {output_filename}")
        print("\n🎉 Your Utah Jazz insights slide is ready!")

        return True


def main():
    """Main entry point"""
    generator = JazzInsightsSlideGenerator()

    # You can customize the output filename here
    output_file = f"Utah_Jazz_Insights_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"

    success = generator.generate_complete_slide(output_file)

    if not success:
        sys.exit(1)


if __name__ == "__main__":
    main()