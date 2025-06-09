"""
Simple test script to verify python-pptx is working correctly
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create a simple presentation
prs = Presentation()

# Add title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
title.text = "Utah Jazz Sponsorship Insights"

subtitle = slide.placeholders[1]
subtitle.text = "Test Presentation"

# Add a slide with the fan wheel image
slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

# Add title
title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
title_box.text = "Fan Wheel Analysis"
title_box.text_frame.paragraphs[0].font.size = Pt(28)
title_box.text_frame.paragraphs[0].font.bold = True

# Add wheel image if it exists
wheel_path = "professional_fan_wheel.png"
if os.path.exists(wheel_path):
    slide2.shapes.add_picture(
        wheel_path,
        Inches(2.5),  # Center horizontally
        Inches(1.5),
        width=Inches(5)
    )
    print(f"✓ Added wheel image from {wheel_path}")
else:
    print(f"✗ Wheel image not found at {wheel_path}")

# Add some text
text_box = slide2.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1))
text_box.text = "Jazz fans are values-driven entertainment seekers!"
text_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
text_box.text_frame.paragraphs[0].font.size = Pt(18)

# Save the presentation
output_file = "test_presentation.pptx"
prs.save(output_file)
print(f"\n✅ Test presentation saved as {output_file}")
print(f"✅ Created {len(prs.slides)} slides")
print("\nYou can now open the PowerPoint file to verify it works!")